import re
import io
import fitz
import html
import xlrd
import random
import string
import olefile
import pandas as pd
from drf_yasg import openapi
from pptx import Presentation
from datetime import datetime
from docx import Document as Dc
from django.conf import settings
from django.utils import timezone
from openpyxl import load_workbook
from django.core.files.uploadedfile import InMemoryUploadedFile


error_code =settings.ERROR_JSON
def get_client_ip(request):
    x_forwarded_for = request.META.get('HTTP_X_FORWARDED_FOR')
    if x_forwarded_for:
        ip = x_forwarded_for.split(',')[0]
    else:
        ip = request.META.get('REMOTE_ADDR')
    return ip

def get_error_msg(code=0):
    return error_code.get(str(code), "Something went wrong")

def get_now_time():
    return datetime.now(tz=timezone.utc)

def generate_password(length=12):
    characters = string.ascii_letters + string.digits + string.punctuation
    password = ''.join(random.choice(characters) for i in range(length))
    return password

def camel_case_to_snake_case(str):
    s1 = re.sub('(.)([A-Z][a-z]+)', r'\1_\2', str)
    return re.sub('([a-z0-9])([A-Z])', r'\1_\2', s1).lower()

def valid_serializer(serializer, error_code=None):
    if serializer.is_valid():
        return serializer.data
    for x, y in serializer.errors.items():
        if error_code is not None:
            raise Exception(error_code)
        error = str(x) + ' : ' + str(y)
        raise Exception(error)

def generate_token(token_length=24):
    valid_chars = string.ascii_lowercase + string.digits + '-'
    token = random.choice(string.ascii_lowercase + string.digits)
    while len(token) < token_length:
        token += random.choice(valid_chars)
    token += random.choice(string.ascii_lowercase + string.digits)
    return token

def create_table_from_string(string):
    rows = string.split("\n")
    if "|" in rows[0]:
        table=[row.split("|") for row in rows]
    else:
        table = [row.split("\t") for row in rows]
    return table

def extract_text_from_pdf(file_contents):
    with fitz.open(None, file_contents) as doc:
        text = ""
        for page in doc:
            text += page.get_text()
        return text

def extract_text_from_xlsx(file_contents):
    wb = load_workbook(filename=io.BytesIO(file_contents))
    return '\n'.join([str(cell.value) for sheet in wb for row in sheet for cell in row if cell.value is not None])

def extract_text_from_xls(file_contents):
    wb = xlrd.open_workbook(file_contents=file_contents)
    text = ''
    for sheet in wb.sheets():
        for row in range(sheet.nrows):
            for col in range(sheet.ncols):
                cell_value = sheet.cell_value(row, col)
                if cell_value:
                    text += str(cell_value) + '\n'
    return text

def extract_text_from_csv(file_contents):
    csv_file = pd.read_csv(file_contents)
    return csv_file.to_markdown()

def extract_text_from_docx(file_contents):
    doc = Dc(io.BytesIO(file_contents))
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_doc(file_contents):
    ole = olefile.OleFileIO(io.BytesIO(file_contents))
    word_stream = ole.openstream('WordDocument')
    word_bytes = word_stream.read()
    ole.close()

    text_content = None
    try:
        text_content = word_bytes.decode('utf-8')
    except:
        pass
    if not text_content:
        try:
            text_content = word_bytes.decode('latin-1')
        except:
            pass
    if not text_content:
        raise ValueError("Unable to decode Word document stream")
    return text_content

def extract_text_from_pptx(file_contents):
    prs = Presentation(io.BytesIO(file_contents))
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    return text_content

def extract_text_from_file(file_contents, filename):
    if filename.endswith('.pdf'):
        return extract_text_from_pdf(file_contents)
    elif filename.endswith('.xlsx'):
        return extract_text_from_xlsx(file_contents)
    elif filename.endswith('.xls'):
        return extract_text_from_xls(file_contents)
    elif filename.endswith('.docx'):
        return extract_text_from_docx(file_contents)
    elif filename.endswith('.doc'):
        return extract_text_from_doc(file_contents)
    elif filename.endswith('.pptx'):
        return extract_text_from_pptx(file_contents)
    elif filename.endswith('.csv'):
        return extract_text_from_csv(file_contents)
    else:
        # Assume it's a text file
        return file_contents.decode('utf-8')
    
def check_email(email):
    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    if re.match(pattern, email):
        return True
    else:
        return False

def extract_username(email):
    parts = email.split("@")
    if len(parts) == 2:
        username = parts[0]
        return username
    else:
        return ""

def create_swagger_params(name, required=True, type='string',header_type="header",extra={}):
    swagger_type = openapi.TYPE_STRING
    header = openapi.IN_HEADER
    if type == "int":
        swagger_type = openapi.TYPE_INTEGER
    elif type == "bool":
        swagger_type = openapi.TYPE_BOOLEAN
    if header_type=="query":
        header=openapi.IN_QUERY
    return openapi.Parameter(name, header, type=swagger_type, required=required,**extra)

def create_swagger_file_params(name, type):
    if type == "multiple":
        return openapi.Parameter(
            name=name,
            in_=openapi.IN_FORM,
            description='Multiple files upload',
            type=openapi.TYPE_ARRAY,
            required=True,
            items=openapi.Schema(type=openapi.TYPE_FILE),
            style="form",
            explode=True,
        )
    elif type == "single":
        return openapi.Parameter(
            name=name,
            in_=openapi.IN_FORM,
            description='The file to upload.',
            type=openapi.TYPE_FILE,
            required=True
        )

def file_name_changer(file):
    file_content = file.read()
    file_names = file.name.split('.')
    file_name = "_".join(file_names[:-1])+"."+file_names[-1]
    return InMemoryUploadedFile(
       io. BytesIO(file_content),
        "file_field",
        file_name,
        file.content_type,
        file.tell,
        file.file 
    )


token = create_swagger_params(name="token", type="string",header_type="header")
jwtToken = create_swagger_params(name="Authorization", type="string",header_type="header")
multiple_files = create_swagger_file_params(name="files", type="multiple")
single_file = create_swagger_file_params(name="file", type="single")

post_login = [jwtToken]
post_upload = [jwtToken, token]
admin_post_login = [jwtToken]